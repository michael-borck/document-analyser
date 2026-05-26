"""Slide-design metrics for .pptx — additive to the existing prose/readability analysis.

Reads the deck through python-pptx and surfaces *design* signals — slide count,
title coverage, words-per-slide density, image counts, bullet depth, layout
diversity, empty slides, text-overload — that the text-only path (markitdown
→ extracted text) can't see.

Returned as a `SlideDesign` block on the document-analyser response; only
populated when the input is a .pptx file. Composable with the existing
readability analysis on the same deck (which still runs on the extracted text).
"""
from __future__ import annotations

import io
from dataclasses import asdict, dataclass, field
from pathlib import Path

# Empirically: PowerPoint authoring guides commonly cite ~30–50 words/slide for
# audience-facing decks. We flag anything above 80 as "text-overloaded" — past
# this threshold, slides are usually being used as documents, not as visuals.
_TEXT_OVERLOAD_WORDS = 80


@dataclass
class SlideStats:
    index: int
    title: str | None
    layout: str
    word_count: int
    image_count: int
    bullet_depth: int
    is_empty: bool
    is_text_overloaded: bool


@dataclass
class SlideDesign:
    slide_count: int = 0
    title_coverage: float = 0.0
    empty_slides: int = 0
    avg_words_per_slide: float = 0.0
    max_words_per_slide: int = 0
    text_overloaded_slides: int = 0
    avg_images_per_slide: float = 0.0
    total_images: int = 0
    max_bullet_depth: int = 0
    distinct_layouts: int = 0
    layout_names: list[str] = field(default_factory=list)
    per_slide: list[SlideStats] = field(default_factory=list)

    def to_dict(self) -> dict:
        d = asdict(self)
        return d


def analyse_pptx(source: bytes | str | Path) -> SlideDesign:
    """Compute slide-design metrics for a .pptx.

    Args:
        source: raw bytes, or a path to the .pptx on disk.

    Returns:
        SlideDesign populated with per-slide + rollup metrics.

    Raises:
        SlideAnalysisError: if the file can't be opened as a .pptx.
    """
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError as e:  # pragma: no cover — declared dep
        raise SlideAnalysisError(f"python-pptx not installed: {e}") from e

    handle = io.BytesIO(source) if isinstance(source, bytes) else str(source)
    try:
        deck = Presentation(handle)
    except Exception as e:
        raise SlideAnalysisError(f"Could not open as .pptx: {e}") from e

    per_slide: list[SlideStats] = []
    layouts: list[str] = []
    total_images = 0
    total_words = 0
    max_words = 0
    max_bullet_depth = 0
    empty_count = 0
    overload_count = 0
    titled_count = 0

    for idx, slide in enumerate(deck.slides, start=1):
        layout_name = (slide.slide_layout.name or "Unknown").strip() or "Unknown"
        layouts.append(layout_name)

        title_shape = None
        try:
            title_shape = slide.shapes.title
        except Exception:
            # Some layouts genuinely don't have a title placeholder.
            title_shape = None
        title_text = None
        if title_shape is not None and title_shape.has_text_frame:
            t = (title_shape.text_frame.text or "").strip()
            if t:
                title_text = t
                titled_count += 1

        slide_word_count = 0
        slide_image_count = 0
        slide_bullet_depth = 0

        for shape in slide.shapes:
            # Images.
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                slide_image_count += 1
                continue
            # Text content.
            if not shape.has_text_frame:
                continue
            tf = shape.text_frame
            for para in tf.paragraphs:
                text = (para.text or "").strip()
                if text:
                    slide_word_count += len(text.split())
                # paragraph.level is 0 for top-level bullets; mostly 0–4 in practice.
                level = getattr(para, "level", 0) or 0
                if level > slide_bullet_depth:
                    slide_bullet_depth = level

        total_images += slide_image_count
        total_words += slide_word_count
        if slide_word_count > max_words:
            max_words = slide_word_count
        if slide_bullet_depth > max_bullet_depth:
            max_bullet_depth = slide_bullet_depth

        is_empty = (slide_word_count == 0 and slide_image_count == 0 and not title_text)
        is_overloaded = slide_word_count > _TEXT_OVERLOAD_WORDS
        if is_empty:
            empty_count += 1
        if is_overloaded:
            overload_count += 1

        per_slide.append(SlideStats(
            index=idx,
            title=title_text,
            layout=layout_name,
            word_count=slide_word_count,
            image_count=slide_image_count,
            bullet_depth=slide_bullet_depth,
            is_empty=is_empty,
            is_text_overloaded=is_overloaded,
        ))

    n = len(per_slide)
    return SlideDesign(
        slide_count=n,
        title_coverage=round(titled_count / n, 4) if n else 0.0,
        empty_slides=empty_count,
        avg_words_per_slide=round(total_words / n, 2) if n else 0.0,
        max_words_per_slide=max_words,
        text_overloaded_slides=overload_count,
        avg_images_per_slide=round(total_images / n, 4) if n else 0.0,
        total_images=total_images,
        max_bullet_depth=max_bullet_depth,
        distinct_layouts=len(set(layouts)),
        layout_names=sorted(set(layouts)),
        per_slide=per_slide,
    )


class SlideAnalysisError(Exception):
    """Raised when slide-design analysis cannot run."""
