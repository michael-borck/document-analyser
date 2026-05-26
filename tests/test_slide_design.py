"""Tests for the .pptx slide-design analyzer.

Each test builds a synthetic .pptx in-memory via python-pptx itself — no binary
fixtures on disk — then asserts the metrics. Includes a CLI smoke + API smoke
to verify the additive `slide_design` block flows end-to-end.
"""
from __future__ import annotations

import io
import json
import subprocess
import sys
from pathlib import Path

import pytest
from fastapi.testclient import TestClient

from document_analyser.analyzers.slide_design import (
    SlideAnalysisError,
    SlideDesign,
    analyse_pptx,
)


# ── synthetic .pptx factories ────────────────────────────────────────────


def _save_deck(prs, tmp_path: Path, name: str = "deck.pptx") -> Path:
    path = tmp_path / name
    prs.save(path)
    return path


def _basic_deck() -> "Presentation":
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    # Slide 1 — title slide
    title_layout = prs.slide_layouts[0]  # "Title Slide"
    s1 = prs.slides.add_slide(title_layout)
    s1.shapes.title.text = "Deck Demo"
    s1.placeholders[1].text = "by Test Author"

    # Slide 2 — title + content with bullets
    content_layout = prs.slide_layouts[1]  # "Title and Content"
    s2 = prs.slides.add_slide(content_layout)
    s2.shapes.title.text = "Bulleted Content"
    body = s2.placeholders[1].text_frame
    body.text = "First top-level"
    p2 = body.add_paragraph()
    p2.text = "Second top-level"
    p3 = body.add_paragraph()
    p3.text = "Nested once"
    p3.level = 1
    p4 = body.add_paragraph()
    p4.text = "Nested twice"
    p4.level = 2

    # Slide 3 — content-only, very wordy (should be text-overloaded)
    s3 = prs.slides.add_slide(content_layout)
    s3.shapes.title.text = "Wall of Text"
    body3 = s3.placeholders[1].text_frame
    body3.text = " ".join(["word"] * 100)  # 100 words → overloaded (threshold = 80)

    # Slide 4 — blank layout, no content (empty)
    blank_layout = prs.slide_layouts[6]  # "Blank"
    prs.slides.add_slide(blank_layout)

    return prs


# ── analyser tests ───────────────────────────────────────────────────────


class TestSlideDesign:
    def test_basic_metrics(self, tmp_path: Path):
        prs = _basic_deck()
        path = _save_deck(prs, tmp_path)

        sd = analyse_pptx(path)
        assert isinstance(sd, SlideDesign)
        assert sd.slide_count == 4
        # 3 of 4 slides have a title (the blank one doesn't)
        assert sd.title_coverage == 0.75
        # Slide 4 is genuinely empty (no title, no text, no image)
        assert sd.empty_slides == 1
        # Slide 3 has 100 words → overloaded
        assert sd.text_overloaded_slides == 1
        # Max bullet depth is 2 (slide 2)
        assert sd.max_bullet_depth == 2

    def test_bytes_input(self, tmp_path: Path):
        prs = _basic_deck()
        buf = io.BytesIO()
        prs.save(buf)
        sd = analyse_pptx(buf.getvalue())
        assert sd.slide_count == 4

    def test_layout_diversity(self, tmp_path: Path):
        prs = _basic_deck()
        path = _save_deck(prs, tmp_path)
        sd = analyse_pptx(path)
        # We used 3 distinct layouts: Title Slide, Title and Content, Blank.
        assert sd.distinct_layouts == 3
        assert "Blank" in sd.layout_names

    def test_per_slide_breakdown(self, tmp_path: Path):
        prs = _basic_deck()
        path = _save_deck(prs, tmp_path)
        sd = analyse_pptx(path)
        per = {s.index: s for s in sd.per_slide}
        assert per[1].title == "Deck Demo"
        assert per[2].bullet_depth == 2
        assert per[3].is_text_overloaded is True
        assert per[4].is_empty is True

    def test_invalid_file_raises(self, tmp_path: Path):
        bogus = tmp_path / "notapptx.pptx"
        bogus.write_bytes(b"not a real pptx")
        with pytest.raises(SlideAnalysisError, match="Could not open"):
            analyse_pptx(bogus)


# ── CLI integration ──────────────────────────────────────────────────────


class TestCLI:
    def test_cli_emits_slide_design_block(self, tmp_path: Path):
        prs = _basic_deck()
        path = _save_deck(prs, tmp_path)
        proc = subprocess.run(
            [sys.executable, "-m", "document_analyser.cli", str(path), "--json"],
            capture_output=True,
            text=True,
        )
        assert proc.returncode == 0, proc.stderr
        data = json.loads(proc.stdout)
        assert "slide_design" in data
        sd = data["slide_design"]
        assert sd["slide_count"] == 4
        assert sd["text_overloaded_slides"] == 1
        assert sd["max_bullet_depth"] == 2


# ── API integration ──────────────────────────────────────────────────────


class TestAPI:
    def test_analyse_route_includes_slide_design(self, tmp_path: Path):
        from document_analyser.api import app

        client = TestClient(app)
        prs = _basic_deck()
        path = _save_deck(prs, tmp_path)

        r = client.post(
            "/analyse",
            files={"file": (path.name, path.read_bytes(),
                            "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
        )
        assert r.status_code == 200, r.text
        body = r.json()
        assert "slide_design" in body
        assert body["slide_design"]["slide_count"] == 4
        # Existing prose/readability still flows.
        assert "word_count" in body
        assert "readability" in body
